from __future__ import annotations

from pathlib import Path
import struct
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from app.services.deck_writer import DeckWriter
from app.services.models import (
    PagePlacementResult,
    PlacementBundle,
    PlacementStatus,
    QcFindingSeverity,
    QcFindingType,
    QcReport,
    SlideQcFinding,
    SlideQcResult,
    SlideQcStatus,
)


def _make_png(path: Path, size: tuple[int, int] = (1200, 700), color: tuple[int, int, int] = (255, 255, 255)) -> Path:
    image = Image.new("RGB", size, color)
    image.save(path)
    return path


def _corrupt_zip_entry(path: Path, entry_name: str) -> None:
    central_directory_struct = struct.Struct("<IHHHHHHIIIHHHHHII")
    payload = bytearray(path.read_bytes())
    entry_name_bytes = entry_name.encode("utf-8")
    offset = 0

    while offset < len(payload):
        signature_offset = payload.find(b"PK\x01\x02", offset)
        if signature_offset == -1:
            break
        header = central_directory_struct.unpack_from(payload, signature_offset)
        file_name_length = header[10]
        extra_field_length = header[11]
        file_comment_length = header[12]
        name_start = signature_offset + central_directory_struct.size
        name_end = name_start + file_name_length
        file_name = bytes(payload[name_start:name_end])
        if file_name == entry_name_bytes:
            crc_offset = signature_offset + 16
            payload[crc_offset] ^= 0x01
            path.write_bytes(payload)
            return
        offset = name_end + extra_field_length + file_comment_length

    raise AssertionError(f"Could not find central directory entry for {entry_name!r}")


def test_build_output_preserves_slides_inserts_pdf_reference_slides_and_appends_unmatched_pdf_pages(tmp_path: Path) -> None:
    source_pptx = tmp_path / "source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, presentation.slide_width, 600000).text_frame.text = "Original slide 1"
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, presentation.slide_width, 600000).text_frame.text = "Original slide 2"
    presentation.save(source_pptx)

    background_image = _make_png(tmp_path / "background-1.png", size=(1600, 900), color=(250, 248, 244))
    second_background_image = _make_png(tmp_path / "background-2.png", size=(1600, 900), color=(245, 242, 236))
    extra_image = _make_png(tmp_path / "extra.png", size=(1600, 900), color=(230, 230, 230))
    writer = DeckWriter()
    bundle = PlacementBundle(
        slide_results=[
            PagePlacementResult(
                candidate_slide_index=0,
                reference_page_index=0,
                status=PlacementStatus.PLACED,
                background_image_path=background_image,
                message="Inserted PDF page as a full-slide reference slide after the original slide",
            ),
            PagePlacementResult(
                candidate_slide_index=1,
                reference_page_index=1,
                status=PlacementStatus.PLACED,
                background_image_path=second_background_image,
                message="Inserted PDF page as a full-slide reference slide after the original slide",
            )
        ],
        extra_reference_results=[
            PagePlacementResult(
                candidate_slide_index=None,
                reference_page_index=2,
                status=PlacementStatus.EXTRA_PDF_PAGE,
                background_image_path=extra_image,
                message="Appended unmatched PDF page as a new reference slide",
            )
        ],
    )

    output_pptx = writer.build_output(source_pptx, bundle, tmp_path / "output.pptx")

    updated = Presentation(output_pptx)
    assert updated.slide_width == presentation.slide_width
    assert len(updated.slides) == 5
    assert len(updated.slides[0].shapes) == 1
    assert len(updated.slides[1].shapes) == 1
    first_reference_slide = updated.slides[1].shapes[0]
    assert first_reference_slide.left == 0
    assert first_reference_slide.top == 0
    assert first_reference_slide.width == updated.slide_width
    assert first_reference_slide.height == updated.slide_height
    assert first_reference_slide.name == "PDF_ORIGINAL"
    assert updated.slides[1]._element.cSld.get("name") == "PDF_ORIGINAL"

    assert len(updated.slides[2].shapes) == 1

    assert len(updated.slides[3].shapes) == 1
    second_reference_slide = updated.slides[3].shapes[0]
    assert second_reference_slide.left == 0
    assert second_reference_slide.top == 0
    assert second_reference_slide.width == updated.slide_width
    assert second_reference_slide.height == updated.slide_height
    assert second_reference_slide.name == "PDF_ORIGINAL"
    assert updated.slides[3]._element.cSld.get("name") == "PDF_ORIGINAL"

    assert len(updated.slides[4].shapes) == 1
    full_slide_picture = updated.slides[4].shapes[0]
    assert full_slide_picture.left == 0
    assert full_slide_picture.top == 0
    assert full_slide_picture.width == updated.slide_width
    assert full_slide_picture.height == updated.slide_height
    assert full_slide_picture.name == "PDF_ORIGINAL"
    assert updated.slides[4]._element.cSld.get("name") == "PDF_ORIGINAL"


def test_build_output_applies_qc_annotations_to_original_slides(tmp_path: Path) -> None:
    source_pptx = tmp_path / "source-qc.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, presentation.slide_width, 600000).text_frame.text = "Original slide 1"
    presentation.save(source_pptx)

    background_image = _make_png(tmp_path / "background-qc.png", size=(1600, 900), color=(250, 248, 244))
    bundle = PlacementBundle(
        slide_results=[
            PagePlacementResult(
                candidate_slide_index=0,
                reference_page_index=0,
                status=PlacementStatus.PLACED,
                background_image_path=background_image,
                message="Inserted PDF page as a full-slide reference slide after the original slide",
            )
        ]
    )
    qc_report = QcReport(
        slide_results=[
            SlideQcResult(
                slide_index=0,
                page_index=0,
                status=SlideQcStatus.FINDINGS,
                findings=[
                    SlideQcFinding(
                        finding_id=1,
                        finding_type=QcFindingType.MISSING_CONTENT,
                        severity=QcFindingSeverity.HIGH,
                        bbox=(0.1, 0.12, 0.3, 0.26),
                        message="Missing block",
                        confidence=0.95,
                    )
                ],
                summary="AI review findings",
                comment_bullets=[
                    "Missing logo in the top-right corner",
                    "Accent bar color does not match the PDF reference",
                ],
                note="AI review findings\n- Missing logo in the top-right corner\n- Accent bar color does not match the PDF reference",
            )
        ],
        counts_by_type={"missing_content": 1},
    )

    output_pptx = DeckWriter().build_output(source_pptx, bundle, tmp_path / "output-qc.pptx", qc_report=qc_report)

    updated = Presentation(output_pptx)
    original_slide = updated.slides[0]
    shape_names = [shape.name for shape in original_slide.shapes]

    assert shape_names.count("PDF_ORIGINAL") == 1
    assert "QC_BADGE_1" not in shape_names
    summary_shape = next(
        shape
        for shape in original_slide.shapes
        if shape.name == "PDF_ORIGINAL" and shape.has_text_frame and any(p.text for p in shape.text_frame.paragraphs)
    )
    paragraph_text = [paragraph.text for paragraph in summary_shape.text_frame.paragraphs if paragraph.text]
    assert paragraph_text == [
        "AI review findings",
        "• Missing logo in the top-right corner",
        "• Accent bar color does not match the PDF reference",
    ]
    assert all(paragraph.alignment == PP_ALIGN.LEFT for paragraph in summary_shape.text_frame.paragraphs if paragraph.text)
    assert summary_shape.left > updated.slide_width * 0.5
    assert summary_shape.top > updated.slide_height * 0.5
    assert len(updated.slides) == 2


def test_build_reference_only_output_by_package_patch_preserves_slides_and_inserts_reference_slides(tmp_path: Path) -> None:
    source_pptx = tmp_path / "source-reference-only.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, presentation.slide_width, 600000).text_frame.text = "Original slide 1"
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, presentation.slide_width, 600000).text_frame.text = "Original slide 2"
    presentation.save(source_pptx)

    first_reference = _make_png(tmp_path / "reference-1.png", size=(1600, 900), color=(250, 248, 244))
    second_reference = _make_png(tmp_path / "reference-2.png", size=(1600, 900), color=(245, 242, 236))
    extra_reference = _make_png(tmp_path / "reference-extra.png", size=(1600, 900), color=(230, 230, 230))

    bundle = PlacementBundle(
        slide_results=[
            PagePlacementResult(
                candidate_slide_index=0,
                reference_page_index=0,
                status=PlacementStatus.PLACED,
                background_image_path=first_reference,
                message="Inserted PDF page after slide 1",
            ),
            PagePlacementResult(
                candidate_slide_index=1,
                reference_page_index=1,
                status=PlacementStatus.PLACED,
                background_image_path=second_reference,
                message="Inserted PDF page after slide 2",
            ),
        ],
        extra_reference_results=[
            PagePlacementResult(
                candidate_slide_index=None,
                reference_page_index=2,
                status=PlacementStatus.EXTRA_PDF_PAGE,
                background_image_path=extra_reference,
                message="Appended unmatched PDF page as a new reference slide",
            )
        ],
    )

    output_pptx = DeckWriter().build_reference_only_output_by_package_patch(
        source_pptx,
        bundle,
        tmp_path / "output-reference-only.pptx",
    )

    updated = Presentation(output_pptx)
    assert len(updated.slides) == 5
    assert updated.slides[1]._element.cSld.get("name") == "PDF_ORIGINAL"
    assert updated.slides[1].shapes[0].name == "PDF_ORIGINAL"
    assert updated.slides[3]._element.cSld.get("name") == "PDF_ORIGINAL"
    assert updated.slides[3].shapes[0].name == "PDF_ORIGINAL"
    assert updated.slides[4]._element.cSld.get("name") == "PDF_ORIGINAL"
    assert updated.slides[4].shapes[0].name == "PDF_ORIGINAL"


def test_build_reference_only_output_by_package_patch_ignores_corrupted_media_entries(tmp_path: Path) -> None:
    source_pptx = tmp_path / "source-corrupted.pptx"
    embedded_image = _make_png(tmp_path / "embedded.png", size=(40, 40), color=(0, 0, 255))

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(embedded_image), 0, 0)
    presentation.save(source_pptx)
    _corrupt_zip_entry(source_pptx, "ppt/media/image1.png")

    with zipfile.ZipFile(source_pptx, "r") as archive:
        assert archive.testzip() == "ppt/media/image1.png"

    reference_image = _make_png(tmp_path / "reference-page.png", size=(1600, 900), color=(250, 248, 244))
    bundle = PlacementBundle(
        slide_results=[
            PagePlacementResult(
                candidate_slide_index=0,
                reference_page_index=0,
                status=PlacementStatus.PLACED,
                background_image_path=reference_image,
                message="Inserted PDF page after slide 1",
            )
        ]
    )

    output_pptx = DeckWriter().build_reference_only_output_by_package_patch(
        source_pptx,
        bundle,
        tmp_path / "output-corrupted-reference-only.pptx",
    )

    with zipfile.ZipFile(output_pptx, "r") as archive:
        presentation_rels_xml = archive.read("ppt/_rels/presentation.xml.rels")
        inserted_slide_xml = archive.read("ppt/slides/slide2.xml")

    assert b'slides/slide2.xml' in presentation_rels_xml
    assert b'name="PDF_ORIGINAL"' in inserted_slide_xml
