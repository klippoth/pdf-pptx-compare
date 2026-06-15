from __future__ import annotations

from pathlib import Path
import re
import zipfile

import pytest
from pptx import Presentation

from app.services.models import PagePlacementResult, PlacementBundle, PlacementStatus
from app.services.renderer import LibreOfficeRenderer, PowerPointRenderer, Renderer, RendererError


class FailingRenderer:
    def is_available(self) -> bool:
        return True

    def export_pptx_to_pdf(self, pptx_path: Path, output_path: Path) -> Path:
        raise RendererError("primary failed")

    def export_pptx_to_slide_images(self, pptx_path: Path, output_dir: Path, *, dpi: int = 180) -> list[Path]:
        raise RendererError("primary failed")


class SuccessfulRenderer:
    def is_available(self) -> bool:
        return True

    def export_pptx_to_pdf(self, pptx_path: Path, output_path: Path) -> Path:
        output_path.write_bytes(b"%PDF-1.4")
        return output_path

    def export_pptx_to_slide_images(self, pptx_path: Path, output_dir: Path, *, dpi: int = 180) -> list[Path]:
        output_dir.mkdir(parents=True, exist_ok=True)
        image_path = output_dir / "Slide1.PNG"
        image_path.write_bytes(b"png")
        return [image_path]

    def build_output_with_reference_slides(self, source_pptx: Path, placement_bundle: PlacementBundle, output_path: Path) -> Path:
        output_path.write_bytes(b"pptx")
        return output_path


class RecordingRenderer:
    def __init__(self, should_fail: bool = False, available: bool = True):
        self.called = False
        self.slide_image_called = False
        self.output_build_called = False
        self.should_fail = should_fail
        self.available = available
        self.mac_slide_export_macro_name = None
        self.mac_reference_slide_insert_macro_name = None

    def is_available(self) -> bool:
        return self.available

    def export_pptx_to_pdf(self, pptx_path: Path, output_path: Path) -> Path:
        self.called = True
        if self.should_fail:
            raise RendererError("renderer failed")
        output_path.write_bytes(b"%PDF-1.4")
        return output_path

    def export_pptx_to_slide_images(self, pptx_path: Path, output_dir: Path, *, dpi: int = 180) -> list[Path]:
        self.slide_image_called = True
        if self.should_fail:
            raise RendererError("renderer failed")
        output_dir.mkdir(parents=True, exist_ok=True)
        image_path = output_dir / "Slide1.PNG"
        image_path.write_bytes(b"png")
        return [image_path]

    def build_output_with_reference_slides(self, source_pptx: Path, placement_bundle: PlacementBundle, output_path: Path) -> Path:
        self.output_build_called = True
        if self.should_fail:
            raise RendererError("renderer failed")
        output_path.write_bytes(b"pptx")
        return output_path


def test_renderer_falls_back_to_secondary(tmp_path: Path) -> None:
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=FailingRenderer(),
        libreoffice=SuccessfulRenderer(),
        enable_powerpoint_fallback=True,
    )
    output_path = tmp_path / "candidate.pdf"

    exported = renderer.export_pptx_to_pdf(tmp_path / "deck.pptx", output_path)

    assert exported == output_path
    assert renderer.last_used_renderer == "libreoffice"


def test_renderer_does_not_call_powerpoint_when_fallback_disabled(tmp_path: Path) -> None:
    powerpoint = RecordingRenderer()
    libreoffice = RecordingRenderer()
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=powerpoint,
        libreoffice=libreoffice,
        enable_powerpoint_fallback=False,
    )

    exported = renderer.export_pptx_to_pdf(tmp_path / "deck.pptx", tmp_path / "candidate.pdf")

    assert exported.exists()
    assert libreoffice.called is True
    assert powerpoint.called is False
    assert renderer.last_used_renderer == "libreoffice"


def test_renderer_uses_powerpoint_when_libreoffice_is_missing(tmp_path: Path) -> None:
    powerpoint = RecordingRenderer(available=True)
    libreoffice = RecordingRenderer(available=False)
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=powerpoint,
        libreoffice=libreoffice,
        enable_powerpoint_fallback=False,
    )

    exported = renderer.export_pptx_to_pdf(tmp_path / "deck.pptx", tmp_path / "candidate.pdf")

    assert exported.exists()
    assert libreoffice.called is False
    assert powerpoint.called is True
    assert renderer.last_used_renderer == "powerpoint"
    assert renderer.availability().can_export_slide_images is False


def test_renderer_can_prefer_powerpoint_without_fallback(tmp_path: Path) -> None:
    powerpoint = RecordingRenderer(available=True)
    libreoffice = RecordingRenderer(available=True)
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=powerpoint,
        libreoffice=libreoffice,
        enable_powerpoint_fallback=False,
    )

    exported = renderer.export_pptx_to_pdf(
        tmp_path / "deck.pptx",
        tmp_path / "candidate.pdf",
        preferred_renderer="powerpoint",
        allow_fallback=False,
    )

    assert exported.exists()
    assert powerpoint.called is True
    assert libreoffice.called is False
    assert renderer.last_used_renderer == "powerpoint"


def test_renderer_rejects_unavailable_preferred_renderer(tmp_path: Path) -> None:
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=RecordingRenderer(available=False),
        libreoffice=RecordingRenderer(available=True),
        enable_powerpoint_fallback=False,
    )

    with pytest.raises(RendererError, match="requested renderer `powerpoint`"):
        renderer.export_pptx_to_pdf(
            tmp_path / "deck.pptx",
            tmp_path / "candidate.pdf",
            preferred_renderer="powerpoint",
            allow_fallback=False,
        )


def test_renderer_reports_install_guidance_when_no_converter_is_available(tmp_path: Path) -> None:
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=RecordingRenderer(available=False),
        libreoffice=RecordingRenderer(available=False),
        enable_powerpoint_fallback=False,
    )

    availability = renderer.availability()

    assert availability.can_convert is False
    assert availability.preferred_renderer == "none"
    assert availability.can_export_slide_images is False
    assert "Install LibreOffice or Microsoft PowerPoint" in availability.message

    with pytest.raises(RendererError, match="Install LibreOffice or Microsoft PowerPoint"):
        renderer.export_pptx_to_pdf(tmp_path / "deck.pptx", tmp_path / "candidate.pdf")


def test_powerpoint_renderer_uses_documented_windows_pdf_export_flow() -> None:
    renderer = PowerPointRenderer(platform_name="windows")

    script = renderer._build_powershell_script(Path(r"C:\deck.pptx"), Path(r"C:\deck.pdf"))

    assert "Presentations.Open($inputPath, 0, 0, 0)" in script
    assert "SaveAs($outputPath, 32)" in script


def test_powerpoint_renderer_uses_documented_windows_slide_image_export_flow() -> None:
    renderer = PowerPointRenderer(platform_name="windows")

    script = renderer._build_powershell_slide_export_script(Path(r"C:\deck.pptx"), Path(r"C:\slides"), dpi=180)

    assert "Presentations.Open($inputPath, 0, 0, 0)" in script
    assert "$presentation.Export($outputDir, 'PNG', $scaleWidth, $scaleHeight)" in script
    assert "$presentation.PageSetup.SlideWidth" in script
    assert "$presentation.PageSetup.SlideHeight" in script


def test_powerpoint_renderer_targets_staged_presentation_for_macro_slide_export(tmp_path: Path) -> None:
    renderer = PowerPointRenderer(
        platform_name="darwin",
        mac_slide_export_macro_name="ExportSlidesToFolder",
    )

    script = renderer._build_applescript_macro_slide_export(
        tmp_path / "codex-123-candidate.pptx",
        tmp_path / "slides",
    )

    assert 'set expectedPresentationName to "codex-123-candidate.pptx"' in script
    assert "set currentPresentation to active presentation" in script
    assert 'if name of currentPresentation is not expectedPresentationName then' in script


def test_powerpoint_renderer_natural_path_sort_handles_double_digits() -> None:
    paths = [
        Path("Slide_1.png"),
        Path("Slide_10.png"),
        Path("Slide_11.png"),
        Path("Slide_2.png"),
    ]

    ordered = sorted(paths, key=PowerPointRenderer._natural_path_sort_key)

    assert [path.name for path in ordered] == [
        "Slide_1.png",
        "Slide_2.png",
        "Slide_10.png",
        "Slide_11.png",
    ]


def test_powerpoint_renderer_validates_exported_slide_count(tmp_path: Path) -> None:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.slides.add_slide(presentation.slide_layouts[6])
    pptx_path = tmp_path / "deck.pptx"
    presentation.save(pptx_path)

    renderer = PowerPointRenderer(platform_name="darwin")
    exported = [tmp_path / "Slide_1.png"]
    exported[0].write_bytes(b"png")

    with pytest.raises(RendererError, match="wrong number of PNG files"):
        renderer._validate_exported_slide_images(pptx_path, exported)


def test_renderer_exports_slide_images_via_powerpoint_only(tmp_path: Path) -> None:
    powerpoint = RecordingRenderer(available=True)
    libreoffice = RecordingRenderer(available=True)
    renderer = Renderer(
        platform_name="windows",
        powerpoint=powerpoint,
        libreoffice=libreoffice,
        enable_powerpoint_fallback=False,
    )

    exported = renderer.export_pptx_to_slide_images(tmp_path / "deck.pptx", tmp_path / "slide-images")

    assert [path.name for path in exported] == ["Slide1.PNG"]
    assert powerpoint.slide_image_called is True
    assert libreoffice.slide_image_called is False


def test_renderer_reports_windows_slide_image_export_availability() -> None:
    renderer = Renderer(
        platform_name="windows",
        powerpoint=RecordingRenderer(available=True),
        libreoffice=RecordingRenderer(available=True),
        enable_powerpoint_fallback=True,
    )

    availability = renderer.availability()

    assert availability.can_export_slide_images is True
    assert availability.slide_image_export_renderer == "powerpoint"
    assert "PowerPoint automation" in availability.slide_image_export_message


def test_renderer_reports_macos_slide_image_export_fallback_message() -> None:
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=RecordingRenderer(available=True),
        libreoffice=RecordingRenderer(available=True),
        enable_powerpoint_fallback=False,
    )

    availability = renderer.availability()

    assert availability.can_export_slide_images is False
    assert "requires a loaded VBA macro add-in" in availability.slide_image_export_message


def test_renderer_reports_macos_macro_slide_image_export_availability() -> None:
    powerpoint = RecordingRenderer(available=True)
    powerpoint.mac_slide_export_macro_name = "ExportSlidesToFolder"
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=powerpoint,
        libreoffice=RecordingRenderer(available=True),
        enable_powerpoint_fallback=False,
    )

    availability = renderer.availability()

    assert availability.can_export_slide_images is True
    assert availability.slide_image_export_renderer == "powerpoint"
    assert "ExportSlidesToFolder" in availability.slide_image_export_message


def test_powerpoint_renderer_builds_macos_macro_slide_export_script(tmp_path: Path) -> None:
    renderer = PowerPointRenderer(
        platform_name="darwin",
        mac_slide_export_macro_name="ExportSlidesToFolder",
        mac_slide_export_long_edge=2800,
    )

    script = renderer._build_applescript_macro_slide_export(tmp_path / "deck.pptx", tmp_path / "slides")

    assert 'set preferredLongEdge to 2800' in script
    assert 'run VB macro macro name "ExportSlidesToFolder" list of parameters {outputFolder, preferredLongEdge}' in script
    assert 'run VB macro macro name "ExportSlidesToFolder" list of parameters {outputFolder}' in script
    assert 'set outputFolder to "' in script
    assert 'mkdir -p ' in script
    assert "open inputPath" not in script

def test_powerpoint_renderer_builds_macos_reference_slide_insert_script(tmp_path: Path) -> None:
    renderer = PowerPointRenderer(
        platform_name="darwin",
        mac_reference_slide_insert_macro_name="InsertReferenceSlidesFromManifest",
    )

    script = renderer._build_applescript_reference_slide_insert(
        tmp_path / "deck-with-pdf-pages.pptx",
        tmp_path / "manifest.tsv",
    )

    assert 'run VB macro macro name "InsertReferenceSlidesFromManifest" list of parameters {manifestPath}' in script
    assert 'set expectedPresentationName to "deck-with-pdf-pages.pptx"' in script
    assert "save currentPresentation" in script


def test_powerpoint_renderer_builds_windows_reference_slide_insert_script(tmp_path: Path) -> None:
    renderer = PowerPointRenderer(platform_name="windows")

    script = renderer._build_powershell_reference_slide_insert_script(
        tmp_path / "source.pptx",
        tmp_path / "output.pptx",
        tmp_path / "manifest.json",
    )

    assert "Copy-Item -LiteralPath $inputPath -Destination $outputPath -Force" in script
    assert "$slide = $presentation.Slides.Add($slideIndex, 12)" in script
    assert "$picture = $slide.Shapes.AddPicture($action.imagePath, 0, -1, 0, 0, $slideWidth, $slideHeight)" in script
    assert "$presentation.Save()" in script


def test_powerpoint_renderer_writes_reference_slide_manifest(tmp_path: Path) -> None:
    image_a = tmp_path / "a.png"
    image_b = tmp_path / "b.png"
    image_a.write_bytes(b"a")
    image_b.write_bytes(b"b")
    actions = [
        PagePlacementResult(
            candidate_slide_index=1,
            reference_page_index=1,
            status=PlacementStatus.PLACED,
            background_image_path=image_b,
            message="insert",
        ),
        PagePlacementResult(
            candidate_slide_index=0,
            reference_page_index=0,
            status=PlacementStatus.PLACED,
            background_image_path=image_a,
            message="insert",
        ),
    ]
    bundle = PlacementBundle(slide_results=actions)

    manifest = PowerPointRenderer._write_reference_slide_manifest(
        PowerPointRenderer._reference_slide_actions(bundle),
        tmp_path / "manifest.tsv",
        tmp_path / "assets",
    )

    lines = manifest.read_text(encoding="utf-8").splitlines()
    assert lines[0].startswith("INSERT\t1\t")
    assert lines[1].startswith("INSERT\t0\t")
    assert (tmp_path / "assets" / "insert-001.png").exists()
    assert (tmp_path / "assets" / "insert-002.png").exists()


def test_renderer_can_build_output_with_reference_slides_on_windows() -> None:
    renderer = Renderer(
        platform_name="windows",
        powerpoint=RecordingRenderer(available=True),
        libreoffice=RecordingRenderer(available=True),
        enable_powerpoint_fallback=False,
    )

    assert renderer.can_build_output_with_reference_slides() is True


def test_renderer_can_build_output_with_reference_slides_on_macos_when_macro_is_configured() -> None:
    powerpoint = RecordingRenderer(available=True)
    powerpoint.mac_reference_slide_insert_macro_name = "InsertReferenceSlidesFromManifest"
    renderer = Renderer(
        platform_name="darwin",
        powerpoint=powerpoint,
        libreoffice=RecordingRenderer(available=True),
        enable_powerpoint_fallback=False,
    )

    assert renderer.can_build_output_with_reference_slides() is True


def test_renderer_build_output_with_reference_slides_delegates_to_powerpoint(tmp_path: Path) -> None:
    powerpoint = RecordingRenderer(available=True)
    renderer = Renderer(
        platform_name="windows",
        powerpoint=powerpoint,
        libreoffice=RecordingRenderer(available=True),
        enable_powerpoint_fallback=False,
    )
    image = tmp_path / "reference.png"
    image.write_bytes(b"png")
    bundle = PlacementBundle(
        slide_results=[
            PagePlacementResult(
                candidate_slide_index=0,
                reference_page_index=0,
                status=PlacementStatus.PLACED,
                background_image_path=image,
                message="insert",
            )
        ]
    )

    output = renderer.build_output_with_reference_slides(tmp_path / "source.pptx", bundle, tmp_path / "output.pptx")

    assert output == tmp_path / "output.pptx"
    assert output.exists()
    assert powerpoint.output_build_called is True


def test_renderer_normalizes_native_reference_slide_names_in_final_pptx(tmp_path: Path) -> None:
    presentation = Presentation()
    while len(presentation.slides) < 2:
        presentation.slides.add_slide(presentation.slide_layouts[6])
    output_path = tmp_path / "native-output.pptx"
    presentation.save(output_path)

    with zipfile.ZipFile(output_path, "r") as source_zip:
        entries = {item.filename: source_zip.read(item.filename) for item in source_zip.infolist()}
        infos = {item.filename: item for item in source_zip.infolist()}

    entries["ppt/slides/slide1.xml"] = re.sub(
        rb'(<(?:[\w]+:)?cSld\b)',
        rb'\1 name="PDF_ORIGINAL_001"',
        entries["ppt/slides/slide1.xml"],
        count=1,
    )
    entries["ppt/slides/slide2.xml"] = re.sub(
        rb'(<(?:[\w]+:)?cSld\b)',
        rb'\1 name="KeepMe"',
        entries["ppt/slides/slide2.xml"],
        count=1,
    )
    with zipfile.ZipFile(output_path, "w") as target_zip:
        for filename, data in entries.items():
            target_zip.writestr(infos[filename], data)

    Renderer._normalize_reference_slide_names(output_path)

    with zipfile.ZipFile(output_path, "r") as final_zip:
        slide1 = final_zip.read("ppt/slides/slide1.xml")
        slide2 = final_zip.read("ppt/slides/slide2.xml")

    assert b'name="PDF_ORIGINAL"' in slide1
    assert b"PDF_ORIGINAL_001" not in slide1
    assert b'name="KeepMe"' in slide2


def test_powerpoint_renderer_finalizes_macos_slide_exports_into_job_folder(tmp_path: Path) -> None:
    staging_dir = tmp_path / "staging"
    staging_dir.mkdir(parents=True)
    exported = []
    for name in ("Slide_1.png", "Slide_2.png"):
        path = staging_dir / name
        path.write_bytes(b"png")
        exported.append(path)

    finalized = PowerPointRenderer._finalize_macos_slide_exports(exported, tmp_path / "job-output")

    assert [path.name for path in finalized] == ["Slide_1.png", "Slide_2.png"]
    assert all(path.exists() for path in finalized)
    assert finalized[0].parent.name == "slides"


def test_powerpoint_renderer_stages_macos_input_inside_powerpoint_container(tmp_path: Path, monkeypatch) -> None:
    sandbox_dir = tmp_path / "sandbox-documents"
    monkeypatch.setattr(
        PowerPointRenderer,
        "_macos_sandbox_documents_dir",
        staticmethod(lambda: sandbox_dir),
    )
    pptx_path = tmp_path / "deck.pptx"
    pptx_path.write_bytes(b"pptx")
    renderer = PowerPointRenderer(platform_name="darwin")

    staged = renderer._stage_macos_input(pptx_path)

    assert staged.parent == sandbox_dir
    assert staged.read_bytes() == b"pptx"
    assert staged.name.endswith("-deck.pptx")

    renderer._cleanup_staged_macos_input(staged)
    assert staged.exists() is False


def test_powerpoint_renderer_rejects_unsupported_platform(tmp_path: Path) -> None:
    renderer = PowerPointRenderer(platform_name="linux")

    with pytest.raises(RendererError, match="not supported"):
        renderer.export_pptx_to_pdf(tmp_path / "deck.pptx", tmp_path / "deck.pdf")


def test_libreoffice_renderer_uses_path_lookup_when_no_explicit_binary(monkeypatch, tmp_path: Path) -> None:
    expected = tmp_path / "soffice.exe"
    expected.write_text("", encoding="utf-8")
    monkeypatch.setattr("app.services.renderer.shutil.which", lambda name: str(expected) if name == "soffice.exe" else None)

    renderer = LibreOfficeRenderer(None)

    assert renderer._resolve_binary() == expected
