from __future__ import annotations

from io import BytesIO
from pathlib import Path
import struct
import threading
import time
import zipfile

import numpy as np
from pptx import Presentation

from fastapi import UploadFile

from app.config import Settings
from app.services.jobs import JobManager
from app.services.models import JobRecord, JobState, PageImage, PlacementStatus, SlideQcResult, SlideQcStatus


def _corrupt_zip_entry_crc(path: Path, entry_name: str) -> None:
    central_directory_struct = struct.Struct("<IHHHHHHIIIHHHHHII")
    payload = bytearray(path.read_bytes())
    entry_name_bytes = entry_name.encode("utf-8")
    offset = 0

    while offset < len(payload):
        signature_offset = payload.find(b"PK\x01\x02", offset)
        if signature_offset == -1:
            break
        header = central_directory_struct.unpack_from(payload, signature_offset)
        file_name_length = header[10]
        extra_field_length = header[11]
        file_comment_length = header[12]
        name_start = signature_offset + central_directory_struct.size
        name_end = name_start + file_name_length
        if bytes(payload[name_start:name_end]) == entry_name_bytes:
            crc_offset = signature_offset + 16
            payload[crc_offset] ^= 0x01
            path.write_bytes(payload)
            return
        offset = name_end + extra_field_length + file_comment_length

    raise AssertionError(f"Could not find central directory entry for {entry_name!r}")


def test_create_job_persists_ai_qc_toggle(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
    )
    manager = JobManager(settings=settings)

    pdf_upload = UploadFile(BytesIO(b"%PDF-1.4"), filename="reference.pdf")
    pptx_upload = UploadFile(BytesIO(b"pptx"), filename="candidate.pptx")

    job = manager.create_job(pdf_upload=pdf_upload, pptx_upload=pptx_upload, enable_ai_qc=False)

    assert job.enable_ai_qc is False
    assert job.input_pdf_path.read_bytes() == b"%PDF-1.4"
    assert job.input_pptx_path.read_bytes() == b"pptx"


def test_create_job_defaults_ai_qc_to_disabled(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
    )
    manager = JobManager(settings=settings)

    pdf_upload = UploadFile(BytesIO(b"%PDF-1.4"), filename="reference.pdf")
    pptx_upload = UploadFile(BytesIO(b"pptx"), filename="candidate.pptx")

    job = manager.create_job(pdf_upload=pdf_upload, pptx_upload=pptx_upload)

    assert job.enable_ai_qc is False


def test_create_job_persists_prompt_override(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
    )
    manager = JobManager(settings=settings)

    pdf_upload = UploadFile(BytesIO(b"%PDF-1.4"), filename="reference.pdf")
    pptx_upload = UploadFile(BytesIO(b"pptx"), filename="candidate.pptx")

    job = manager.create_job(
        pdf_upload=pdf_upload,
        pptx_upload=pptx_upload,
        qc_prompt_override="  Focus especially on superscripts and names.  ",
    )

    assert job.qc_prompt_override == "Focus especially on superscripts and names."


def test_create_job_persists_full_prompt_config(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
    )
    manager = JobManager(settings=settings)

    pdf_upload = UploadFile(BytesIO(b"%PDF-1.4"), filename="reference.pdf")
    pptx_upload = UploadFile(BytesIO(b"pptx"), filename="candidate.pptx")

    job = manager.create_job(
        pdf_upload=pdf_upload,
        pptx_upload=pptx_upload,
        qc_prompt_config={
            "general_system_prompt": "visual system",
            "general_user_prompt": "visual user",
            "text_system_prompt": "text system",
            "text_user_prompt": "text user",
        },
    )

    assert job.qc_prompt_config == {
        "general_system_prompt": "visual system",
        "general_user_prompt": "visual user",
        "text_system_prompt": "text system",
        "text_user_prompt": "text user",
    }


def test_run_openai_qc_uses_fresh_evaluator_instances_in_parallel(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
        openai_api_key="test-key",
        openai_qc_parallelism=4,
    )
    manager = JobManager(settings=settings)
    manager._update_job = lambda *args, **kwargs: None

    created_instance_ids: list[int] = []
    active_calls = 0
    max_active_calls = 0
    active_lock = threading.Lock()

    class FakeEvaluator:
        def __init__(self, instance_id: int):
            self.instance_id = instance_id

        def compare_pages(
            self,
            *,
            slide_index,
            page_index,
            reference_page,
            candidate_page,
            candidate_layout=None,
            debug_output_dir,
            prompt_override=None,
            prompt_config=None,
        ):
            nonlocal active_calls, max_active_calls
            with active_lock:
                active_calls += 1
                max_active_calls = max(max_active_calls, active_calls)
            time.sleep(0.05)
            with active_lock:
                active_calls -= 1
            return SlideQcResult(
                slide_index=slide_index,
                page_index=page_index,
                status=SlideQcStatus.OK,
                note=f"instance-{self.instance_id}",
            )

    def build_fake_evaluator():
        instance_id = len(created_instance_ids) + 1
        created_instance_ids.append(instance_id)
        return FakeEvaluator(instance_id)

    manager._build_openai_qc_evaluator = build_fake_evaluator

    def page(i: int) -> PageImage:
        image = np.full((10, 10, 3), 255, dtype=np.uint8)
        return PageImage(page_index=i, image=image, image_path=tmp_path / f"page-{i}.png")

    matched_pairs = [(i, page(i), page(i)) for i in range(4)]

    results = manager._run_openai_qc(
        "job-123",
        matched_pairs,
        candidate_layouts=[],
        prompt_override="Focus on titles",
        prompt_config={"text_user_prompt": "Review lines one by one"},
    )

    assert len(created_instance_ids) == 4
    assert {result.note for result in results.values()} == {"instance-1", "instance-2", "instance-3", "instance-4"}
    assert max_active_calls > 1


def test_publish_output_pptx_copies_to_downloads_with_preserved_filename(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
    )
    manager = JobManager(settings=settings)
    built_output = tmp_path / "runs" / "job-1" / "output" / "candidate_with_pdf_pages.pptx"
    built_output.parent.mkdir(parents=True, exist_ok=True)
    built_output.write_bytes(b"first-version")

    published = manager._publish_output_pptx(built_output, "candidate_with_pdf_pages.pptx")

    assert published == settings.downloads_dir / "candidate_with_pdf_pages.pptx"
    assert published.read_bytes() == b"first-version"


def test_publish_output_pptx_overwrites_existing_download_copy(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
    )
    manager = JobManager(settings=settings)
    built_output = tmp_path / "runs" / "job-1" / "output" / "candidate_with_pdf_pages.pptx"
    built_output.parent.mkdir(parents=True, exist_ok=True)
    built_output.write_bytes(b"updated-version")
    existing_download = settings.downloads_dir / "candidate_with_pdf_pages.pptx"
    existing_download.parent.mkdir(parents=True, exist_ok=True)
    existing_download.write_bytes(b"old-version")

    published = manager._publish_output_pptx(built_output, "candidate_with_pdf_pages.pptx")

    assert published.read_bytes() == b"updated-version"


def test_build_reference_only_placement_bundle_uses_slide_metrics_without_candidate_rendering(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
        render_dpi=180,
    )
    manager = JobManager(settings=settings)

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.slides.add_slide(presentation.slide_layouts[6])
    pptx_path = tmp_path / "candidate.pptx"
    presentation.save(pptx_path)

    reference_pages = []
    for index in range(3):
        image = np.full((120, 200, 3), 255 - (index * 10), dtype=np.uint8)
        reference_pages.append(PageImage(page_index=index, image=image, image_path=tmp_path / f"reference-{index}.png"))

    bundle, slide_qc_results_by_index, total_pages = manager._build_reference_only_placement_bundle(
        pptx_path,
        reference_pages,
        tmp_path / "placed",
    )

    assert total_pages == 3
    assert len(bundle.slide_results) == 2
    assert len(bundle.extra_reference_results) == 1
    assert slide_qc_results_by_index == {}
    assert all(result.status == PlacementStatus.PLACED for result in bundle.slide_results)
    assert all(result.background_image_path and result.background_image_path.exists() for result in bundle.slide_results)
    assert bundle.extra_reference_results[0].status == PlacementStatus.EXTRA_PDF_PAGE
    assert bundle.extra_reference_results[0].background_image_path.exists()


def test_build_reference_only_placement_bundle_reads_metrics_without_touching_corrupted_media(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
        render_dpi=180,
    )
    manager = JobManager(settings=settings)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    embedded_image_path = tmp_path / "embedded.png"
    from PIL import Image
    Image.new("RGB", (40, 40), (0, 0, 255)).save(embedded_image_path)
    slide.shapes.add_picture(str(embedded_image_path), 0, 0)
    pptx_path = tmp_path / "candidate-corrupted.pptx"
    presentation.save(pptx_path)
    _corrupt_zip_entry_crc(pptx_path, "ppt/media/image1.png")

    with zipfile.ZipFile(pptx_path, "r") as archive:
        assert archive.testzip() == "ppt/media/image1.png"

    reference_image = np.full((120, 200, 3), 255, dtype=np.uint8)
    reference_pages = [PageImage(page_index=0, image=reference_image, image_path=tmp_path / "reference-0.png")]

    bundle, slide_qc_results_by_index, total_pages = manager._build_reference_only_placement_bundle(
        pptx_path,
        reference_pages,
        tmp_path / "placed-corrupted",
    )

    assert total_pages == 1
    assert len(bundle.slide_results) == 1
    assert bundle.slide_results[0].status == PlacementStatus.PLACED
    assert bundle.slide_results[0].background_image_path.exists()
    assert slide_qc_results_by_index == {}


def test_process_job_skips_candidate_rendering_when_ai_qc_is_disabled(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="linux",
    )
    manager = JobManager(settings=settings)

    class FakeRenderer:
        def __init__(self) -> None:
            self.export_pdf_called = False
            self.export_slide_images_called = False
            self.native_output_build_called = False

        def can_export_slide_images(self) -> bool:
            return True

        def export_pptx_to_pdf(self, *args, **kwargs):
            self.export_pdf_called = True
            raise AssertionError("Non-AI processing should not render the candidate deck to PDF.")

        def export_pptx_to_slide_images(self, *args, **kwargs):
            self.export_slide_images_called = True
            raise AssertionError("Non-AI processing should not export candidate slide images.")

        def can_build_output_with_reference_slides(self) -> bool:
            return True

        def build_output_with_reference_slides(self, *args, **kwargs):
            self.native_output_build_called = True
            raise AssertionError("Non-AI processing should not use native PowerPoint output assembly.")

        def is_powerpoint_available(self) -> bool:
            return False

    class FakeRasterizer:
        def __init__(self, pages: list[PageImage]) -> None:
            self.pages = pages
            self.last_used_engine = "fitz"

        def render_pdf(self, *args, **kwargs) -> list[PageImage]:
            return list(self.pages)

    class FakeDeckWriter:
        def __init__(self) -> None:
            self.called = False
            self.reference_only_called = False

        def can_build_reference_only_output_by_package_patch(self) -> bool:
            return True

        def build_reference_only_output_by_package_patch(self, source_pptx: Path, bundle, output_path: Path) -> None:
            self.reference_only_called = True
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(b"pptx")

        def build_output(self, source_pptx: Path, bundle, output_path: Path, qc_report=None) -> None:
            self.called = True
            raise AssertionError("Non-AI processing should prefer the reference-only package patch writer.")

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.slides.add_slide(presentation.slide_layouts[6])

    working_dir = settings.runs_dir / "job-1"
    input_dir = working_dir / "input"
    output_dir = working_dir / "output"
    input_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = input_dir / "reference.pdf"
    pptx_path = input_dir / "candidate.pptx"
    pdf_path.write_bytes(b"%PDF-1.4")
    presentation.save(pptx_path)

    reference_pages = []
    for index in range(2):
        image = np.full((180, 320, 3), 255 - (index * 10), dtype=np.uint8)
        reference_pages.append(PageImage(page_index=index, image=image, image_path=working_dir / f"reference-{index}.png"))

    job = JobRecord(
        job_id="job-1",
        working_dir=working_dir,
        input_pdf_path=pdf_path,
        input_pptx_path=pptx_path,
        original_pptx_name="candidate.pptx",
        enable_ai_qc=False,
    )
    manager._jobs[job.job_id] = job
    manager.renderer = FakeRenderer()
    manager.rasterizer = FakeRasterizer(reference_pages)
    manager.deck_writer = FakeDeckWriter()
    manager._open_output_pptx = lambda output_path: None

    manager._process_job(job.job_id)

    assert manager.renderer.export_pdf_called is False
    assert manager.renderer.export_slide_images_called is False
    assert manager.renderer.native_output_build_called is False
    assert manager.deck_writer.reference_only_called is True
    assert manager.deck_writer.called is False
    assert job.status == JobState.COMPLETED.value
    assert job.output_pptx_path == settings.downloads_dir / "candidate_with_pdf_pages.pptx"
    assert job.output_pptx_path.exists()


def test_process_job_retries_with_native_powerpoint_output_when_crc_error_occurs(tmp_path: Path) -> None:
    settings = Settings(
        base_dir=tmp_path,
        static_dir=tmp_path,
        runs_dir=tmp_path / "runs",
        downloads_dir=tmp_path / "Downloads",
        platform_name="darwin",
    )
    manager = JobManager(settings=settings)

    class FakeRenderer:
        def __init__(self) -> None:
            self.native_output_build_called = False

        def can_export_slide_images(self) -> bool:
            return True

        def can_build_output_with_reference_slides(self) -> bool:
            return True

        def build_output_with_reference_slides(self, source_pptx: Path, bundle, output_path: Path) -> None:
            self.native_output_build_called = True
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(b"native-pptx")

        def is_powerpoint_available(self) -> bool:
            return True

    class FakeRasterizer:
        def __init__(self, pages: list[PageImage]) -> None:
            self.pages = pages
            self.last_used_engine = "fitz"

        def render_pdf(self, *args, **kwargs) -> list[PageImage]:
            return list(self.pages)

    class FakeDeckWriter:
        def __init__(self) -> None:
            self.reference_only_called = False

        def can_build_reference_only_output_by_package_patch(self) -> bool:
            return True

        def build_reference_only_output_by_package_patch(self, source_pptx: Path, bundle, output_path: Path) -> None:
            self.reference_only_called = True
            raise zipfile.BadZipFile("Bad CRC-32 for file 'ppt/media/image48.png'")

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])

    working_dir = settings.runs_dir / "job-2"
    input_dir = working_dir / "input"
    output_dir = working_dir / "output"
    input_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = input_dir / "reference.pdf"
    pptx_path = input_dir / "candidate.pptx"
    pdf_path.write_bytes(b"%PDF-1.4")
    presentation.save(pptx_path)

    reference_image = np.full((180, 320, 3), 255, dtype=np.uint8)
    reference_pages = [PageImage(page_index=0, image=reference_image, image_path=working_dir / "reference-0.png")]

    job = JobRecord(
        job_id="job-2",
        working_dir=working_dir,
        input_pdf_path=pdf_path,
        input_pptx_path=pptx_path,
        original_pptx_name="candidate.pptx",
        enable_ai_qc=False,
    )
    manager._jobs[job.job_id] = job
    manager.renderer = FakeRenderer()
    manager.rasterizer = FakeRasterizer(reference_pages)
    manager.deck_writer = FakeDeckWriter()
    manager._open_output_pptx = lambda output_path: None

    manager._process_job(job.job_id)

    assert manager.deck_writer.reference_only_called is True
    assert manager.renderer.native_output_build_called is True
    assert job.status == JobState.COMPLETED.value
    assert job.output_pptx_path == settings.downloads_dir / "candidate_with_pdf_pages.pptx"
    assert job.output_pptx_path.read_bytes() == b"native-pptx"
