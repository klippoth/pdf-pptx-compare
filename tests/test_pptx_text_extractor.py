from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.services.pptx_text_extractor import PptxTextExtractor


def test_pptx_text_extractor_reads_slide_textboxes_and_tables(tmp_path: Path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(1000000, 800000, 5000000, 1200000)
    textbox.text_frame.text = "Federal Signal Q2 2023 Earnings Call"

    table_shape = slide.shapes.add_table(1, 1, 1000000, 2500000, 4000000, 1000000)
    table_shape.table.cell(0, 0).text = "Jennifer Sherman"
    presentation.save(pptx_path)

    layouts = PptxTextExtractor().extract_document(pptx_path)

    assert len(layouts) == 1
    texts = [paragraph.text for paragraph in layouts[0].paragraphs]
    assert "Federal Signal Q2 2023 Earnings Call" in texts
    assert "Jennifer Sherman" in texts
    assert layouts[0].total_characters > 10
