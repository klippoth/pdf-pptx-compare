from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from app.services.models import SlideQcResult, SlideQcStatus


class AnnotationWriter:
    def __init__(self):
        self.annotation_shape_name = "PDF_ORIGINAL"

    def apply(self, slide, slide_qc_result: SlideQcResult, slide_width, slide_height) -> None:
        if slide_qc_result.status == SlideQcStatus.FINDINGS and (
            slide_qc_result.summary or slide_qc_result.comment_bullets or slide_qc_result.note
        ):
            self._add_summary_note(
                slide,
                slide_qc_result.summary,
                slide_qc_result.comment_bullets,
                slide_qc_result.note,
                slide_width,
                slide_height,
            )
        if slide_qc_result.status == SlideQcStatus.MANUAL_REVIEW and slide_qc_result.note:
            self._add_manual_review_note(slide, slide_qc_result.note, slide_width, slide_height)

    def _add_summary_note(self, slide, summary: str | None, bullets: list[str], note: str | None, slide_width, slide_height) -> None:
        rendered_bullets = [bullet.strip() for bullet in bullets if bullet and bullet.strip()]
        rendered_summary = summary.strip() if summary and summary.strip() else None
        if rendered_summary is None and note:
            rendered_summary = note.strip().splitlines()[0].strip()
        bullet_count = max(1, len(rendered_bullets))
        note_width = int(slide_width * 0.42)
        note_height = int(slide_height * min(0.42, 0.10 + (bullet_count * 0.045)))
        margin_x = int(slide_width * 0.03)
        margin_y = int(slide_height * 0.03)
        left = slide_width - note_width - margin_x
        top = slide_height - note_height - margin_y
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, note_width, note_height)
        box.name = self.annotation_shape_name
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 249, 231)
        box.line.color.rgb = RGBColor(212, 196, 140)

        text_frame = box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.margin_left = Pt(10)
        text_frame.margin_right = Pt(10)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(8)

        if rendered_summary:
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = rendered_summary
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor(58, 58, 58)

        for bullet in rendered_bullets:
            paragraph = text_frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.color.rgb = RGBColor(58, 58, 58)

        if not rendered_bullets and note:
            paragraph = text_frame.add_paragraph() if rendered_summary else text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = note if not rendered_summary else f"• {note}"
            run.font.size = Pt(11 if rendered_summary else 12)
            run.font.bold = False
            run.font.color.rgb = RGBColor(58, 58, 58)

    def _add_manual_review_note(self, slide, note: str, slide_width, slide_height) -> None:
        note_width = int(slide_width * 0.38)
        note_height = int(slide_height * 0.08)
        margin_x = int(slide_width * 0.03)
        margin_y = int(slide_height * 0.03)
        left = slide_width - note_width - margin_x
        top = slide_height - note_height - margin_y
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, note_width, note_height)
        pill.name = self.annotation_shape_name
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(88, 94, 104)
        pill.line.color.rgb = RGBColor(88, 94, 104)
        text_frame = pill.text_frame
        text_frame.clear()
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.margin_left = Pt(10)
        text_frame.margin_right = Pt(10)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(8)
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = note
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
