from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation

from app.services.annotation_writer import AnnotationWriter
from app.services.models import PagePlacementResult, PlacementBundle, PlacementStatus, QcReport


CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
PACKAGE_RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
OFFICE_RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SLIDE_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
SLIDE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
SLIDE_LAYOUT_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
IMAGE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
IMAGE_CONTENT_TYPES = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
}

ET.register_namespace("a", DRAWING_NS)
ET.register_namespace("p", PRESENTATION_NS)
ET.register_namespace("r", OFFICE_RELS_NS)


class DeckWriter:
    def __init__(self):
        self.pdf_reference_shape_name = "PDF_ORIGINAL"
        self.pdf_reference_slide_name = "PDF_ORIGINAL"
        self.annotation_writer = AnnotationWriter()

    def can_build_reference_only_output_by_package_patch(self) -> bool:
        return shutil.which("zip") is not None

    def build_output(
        self,
        source_pptx: Path,
        placement_bundle: PlacementBundle,
        output_path: Path,
        qc_report: QcReport | None = None,
    ) -> Path:
        presentation = Presentation(str(source_pptx))
        original_slides = list(presentation.slides)
        reference_insertions: list[tuple[int, PagePlacementResult]] = []

        for slide_index, slide in enumerate(original_slides):
            if slide_index >= len(placement_bundle.slide_results):
                break
            result = placement_bundle.slide_results[slide_index]
            if qc_report is not None and slide_index < len(qc_report.slide_results):
                self.annotation_writer.apply(
                    slide=slide,
                    slide_qc_result=qc_report.slide_results[slide_index],
                    slide_width=presentation.slide_width,
                    slide_height=presentation.slide_height,
                )
            if result.background_image_path and result.status == PlacementStatus.PLACED:
                reference_insertions.append((slide_index, result))

        for slide_index, result in reversed(reference_insertions):
            self._insert_pdf_reference_slide_after(
                presentation=presentation,
                insert_after_index=slide_index,
                result=result,
            )

        for result in placement_bundle.extra_reference_results:
            self._append_pdf_only_slide(
                presentation=presentation,
                result=result,
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))
        return output_path

    def build_reference_only_output_by_package_patch(
        self,
        source_pptx: Path,
        placement_bundle: PlacementBundle,
        output_path: Path,
    ) -> Path:
        if not self.can_build_reference_only_output_by_package_patch():
            raise RuntimeError("The local zip command is required for reference-only output patching.")

        reference_insertions = [
            (index, result)
            for index, result in enumerate(placement_bundle.slide_results)
            if result.background_image_path and result.status == PlacementStatus.PLACED
        ]
        extra_reference_results = [
            result
            for result in placement_bundle.extra_reference_results
            if result.background_image_path and result.status == PlacementStatus.EXTRA_PDF_PAGE
        ]

        output_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source_pptx, output_path)
        if not reference_insertions and not extra_reference_results:
            return output_path

        with zipfile.ZipFile(source_pptx, "r") as source_zip:
            blank_layout_target = self._find_blank_slide_layout_target(source_zip)
            slide_width, slide_height = self._read_slide_size_emu(source_zip)

            content_types_root = ET.fromstring(source_zip.read("[Content_Types].xml"))
            presentation_root = ET.fromstring(source_zip.read("ppt/presentation.xml"))
            presentation_rels_root = ET.fromstring(source_zip.read("ppt/_rels/presentation.xml.rels"))

            next_slide_number = self._next_part_number(source_zip.namelist(), r"ppt/slides/slide(\d+)\.xml")
            next_media_number = self._next_part_number(source_zip.namelist(), r"ppt/media/image(\d+)\.[^.]+")
            next_relationship_number = self._next_relationship_number(presentation_rels_root)
            next_slide_id = self._next_slide_id(presentation_root)

            slide_id_list = presentation_root.find(f"{{{PRESENTATION_NS}}}sldIdLst")
            if slide_id_list is None:
                raise RuntimeError("The PowerPoint package is missing the slide list in presentation.xml.")

            original_slide_entries = list(slide_id_list)
            patched_slide_entries = []
            staged_entries: list[tuple[str, bytes]] = []

            for slide_index, slide_entry in enumerate(original_slide_entries):
                patched_slide_entries.append(slide_entry)
                if slide_index >= len(placement_bundle.slide_results):
                    continue
                result = placement_bundle.slide_results[slide_index]
                if not (result.background_image_path and result.status == PlacementStatus.PLACED):
                    continue

                slide_part_path = f"ppt/slides/slide{next_slide_number}.xml"
                slide_rels_part_path = f"ppt/slides/_rels/slide{next_slide_number}.xml.rels"
                media_extension = result.background_image_path.suffix.lower().lstrip(".") or "png"
                media_part_path = f"ppt/media/image{next_media_number}.{media_extension}"
                presentation_relationship_id = f"rId{next_relationship_number}"

                patched_slide_entries.append(
                    ET.Element(
                        f"{{{PRESENTATION_NS}}}sldId",
                        {
                            "id": str(next_slide_id),
                            f"{{{OFFICE_RELS_NS}}}id": presentation_relationship_id,
                        },
                    )
                )
                ET.SubElement(
                    presentation_rels_root,
                    f"{{{PACKAGE_RELS_NS}}}Relationship",
                    {
                        "Id": presentation_relationship_id,
                        "Type": SLIDE_REL_TYPE,
                        "Target": f"slides/slide{next_slide_number}.xml",
                    },
                )
                self._ensure_slide_override(content_types_root, slide_part_path)
                self._ensure_image_default(content_types_root, media_extension)
                staged_entries.extend(
                    [
                        (
                            slide_part_path,
                            self._build_reference_slide_xml(
                                slide_width=slide_width,
                                slide_height=slide_height,
                                image_relationship_id="rId2",
                                image_description=result.background_image_path.name,
                            ),
                        ),
                        (
                            slide_rels_part_path,
                            self._build_reference_slide_rels_xml(
                                layout_target=blank_layout_target,
                                image_target=f"../media/{Path(media_part_path).name}",
                            ),
                        ),
                        (media_part_path, result.background_image_path.read_bytes()),
                    ]
                )

                next_slide_number += 1
                next_media_number += 1
                next_relationship_number += 1
                next_slide_id += 1

            for result in extra_reference_results:
                slide_part_path = f"ppt/slides/slide{next_slide_number}.xml"
                slide_rels_part_path = f"ppt/slides/_rels/slide{next_slide_number}.xml.rels"
                media_extension = result.background_image_path.suffix.lower().lstrip(".") or "png"
                media_part_path = f"ppt/media/image{next_media_number}.{media_extension}"
                presentation_relationship_id = f"rId{next_relationship_number}"

                patched_slide_entries.append(
                    ET.Element(
                        f"{{{PRESENTATION_NS}}}sldId",
                        {
                            "id": str(next_slide_id),
                            f"{{{OFFICE_RELS_NS}}}id": presentation_relationship_id,
                        },
                    )
                )
                ET.SubElement(
                    presentation_rels_root,
                    f"{{{PACKAGE_RELS_NS}}}Relationship",
                    {
                        "Id": presentation_relationship_id,
                        "Type": SLIDE_REL_TYPE,
                        "Target": f"slides/slide{next_slide_number}.xml",
                    },
                )
                self._ensure_slide_override(content_types_root, slide_part_path)
                self._ensure_image_default(content_types_root, media_extension)
                staged_entries.extend(
                    [
                        (
                            slide_part_path,
                            self._build_reference_slide_xml(
                                slide_width=slide_width,
                                slide_height=slide_height,
                                image_relationship_id="rId2",
                                image_description=result.background_image_path.name,
                            ),
                        ),
                        (
                            slide_rels_part_path,
                            self._build_reference_slide_rels_xml(
                                layout_target=blank_layout_target,
                                image_target=f"../media/{Path(media_part_path).name}",
                            ),
                        ),
                        (media_part_path, result.background_image_path.read_bytes()),
                    ]
                )

                next_slide_number += 1
                next_media_number += 1
                next_relationship_number += 1
                next_slide_id += 1

            for child in list(slide_id_list):
                slide_id_list.remove(child)
            for child in patched_slide_entries:
                slide_id_list.append(child)

            staged_entries.extend(
                [
                    (
                        "[Content_Types].xml",
                        self._serialize_xml(content_types_root, default_namespace=CONTENT_TYPES_NS),
                    ),
                    ("ppt/presentation.xml", self._serialize_xml(presentation_root)),
                    (
                        "ppt/_rels/presentation.xml.rels",
                        self._serialize_xml(presentation_rels_root, default_namespace=PACKAGE_RELS_NS),
                    ),
                ]
            )

        with tempfile.TemporaryDirectory(prefix="pptx-reference-patch-") as staging_root:
            staging_root_path = Path(staging_root)
            relative_paths: list[str] = []
            for relative_path, payload in staged_entries:
                target_path = staging_root_path / relative_path
                target_path.parent.mkdir(parents=True, exist_ok=True)
                target_path.write_bytes(payload)
                relative_paths.append(relative_path)

            subprocess.run(
                ["zip", "-q", str(output_path), *relative_paths],
                cwd=staging_root,
                check=True,
            )

        return output_path

    def _append_pdf_only_slide(
        self,
        presentation: Presentation,
        result: PagePlacementResult,
    ) -> None:
        slide = self._append_blank_slide(presentation)
        self._set_slide_name(slide, self.pdf_reference_slide_name)
        if result.background_image_path:
            self._add_full_slide_picture(
                slide=slide,
                image_path=result.background_image_path,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
            )

    def _insert_pdf_reference_slide_after(
        self,
        presentation: Presentation,
        insert_after_index: int,
        result: PagePlacementResult,
    ) -> None:
        slide = self._append_blank_slide(presentation)
        self._set_slide_name(slide, self.pdf_reference_slide_name)
        if result.background_image_path:
            self._add_full_slide_picture(
                slide=slide,
                image_path=result.background_image_path,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
            )
        self._move_last_slide_to_index(presentation, insert_after_index + 1)

    def _append_blank_slide(self, presentation: Presentation):
        layout = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]
        return presentation.slides.add_slide(layout)

    def _move_last_slide_to_index(self, presentation: Presentation, insert_index: int) -> None:
        slide_id_list = presentation.slides._sldIdLst
        slide_id = slide_id_list[-1]
        slide_id_list.remove(slide_id)
        slide_id_list.insert(insert_index, slide_id)

    def _add_full_slide_picture(self, slide, image_path: Path, slide_width, slide_height) -> None:
        picture = slide.shapes.add_picture(str(image_path), 0, 0, width=slide_width, height=slide_height)
        picture.name = self.pdf_reference_shape_name

    def _set_slide_name(self, slide, name: str) -> None:
        slide._element.cSld.set("name", name)

    def _find_blank_slide_layout_target(self, source_zip: zipfile.ZipFile) -> str:
        slide_layout_paths = sorted(
            path
            for path in source_zip.namelist()
            if re.fullmatch(r"ppt/slideLayouts/slideLayout\d+\.xml", path)
        )
        if not slide_layout_paths:
            raise RuntimeError("The PowerPoint package has no slide layouts available for reference-slide insertion.")

        fallback_target = f"../slideLayouts/{Path(slide_layout_paths[0]).name}"
        for slide_layout_path in slide_layout_paths:
            root = ET.fromstring(source_zip.read(slide_layout_path))
            if root.get("type") == "blank":
                return f"../slideLayouts/{Path(slide_layout_path).name}"
        return fallback_target

    def _read_slide_size_emu(self, source_zip: zipfile.ZipFile) -> tuple[int, int]:
        presentation_root = ET.fromstring(source_zip.read("ppt/presentation.xml"))
        slide_size = presentation_root.find(f"{{{PRESENTATION_NS}}}sldSz")
        if slide_size is None:
            raise RuntimeError("The PowerPoint package is missing slide size metadata.")
        return int(slide_size.get("cx", "0")), int(slide_size.get("cy", "0"))

    @staticmethod
    def _next_part_number(entry_names: list[str], pattern: str) -> int:
        matcher = re.compile(pattern)
        numbers = [int(match.group(1)) for name in entry_names if (match := matcher.fullmatch(name))]
        return (max(numbers) + 1) if numbers else 1

    @staticmethod
    def _next_relationship_number(presentation_rels_root: ET.Element) -> int:
        numbers = []
        for relationship in presentation_rels_root.findall(f"{{{PACKAGE_RELS_NS}}}Relationship"):
            relationship_id = relationship.get("Id", "")
            match = re.fullmatch(r"rId(\d+)", relationship_id)
            if match:
                numbers.append(int(match.group(1)))
        return (max(numbers) + 1) if numbers else 1

    @staticmethod
    def _next_slide_id(presentation_root: ET.Element) -> int:
        slide_id_list = presentation_root.find(f"{{{PRESENTATION_NS}}}sldIdLst")
        if slide_id_list is None:
            return 256
        slide_ids = [int(entry.get("id", "255")) for entry in slide_id_list.findall(f"{{{PRESENTATION_NS}}}sldId")]
        return (max(slide_ids) + 1) if slide_ids else 256

    def _ensure_slide_override(self, content_types_root: ET.Element, slide_part_path: str) -> None:
        part_name = f"/{slide_part_path}"
        for override in content_types_root.findall(f"{{{CONTENT_TYPES_NS}}}Override"):
            if override.get("PartName") == part_name:
                return
        ET.SubElement(
            content_types_root,
            f"{{{CONTENT_TYPES_NS}}}Override",
            {"PartName": part_name, "ContentType": SLIDE_CONTENT_TYPE},
        )

    def _ensure_image_default(self, content_types_root: ET.Element, extension: str) -> None:
        extension = extension.lower()
        content_type = IMAGE_CONTENT_TYPES.get(extension)
        if content_type is None:
            raise RuntimeError(f"Unsupported reference image format for PowerPoint insertion: .{extension}")
        for default in content_types_root.findall(f"{{{CONTENT_TYPES_NS}}}Default"):
            if default.get("Extension", "").lower() == extension:
                return
        ET.SubElement(
            content_types_root,
            f"{{{CONTENT_TYPES_NS}}}Default",
            {"Extension": extension, "ContentType": content_type},
        )

    def _build_reference_slide_xml(
        self,
        *,
        slide_width: int,
        slide_height: int,
        image_relationship_id: str,
        image_description: str,
    ) -> bytes:
        root = ET.Element(f"{{{PRESENTATION_NS}}}sld")
        common_slide = ET.SubElement(root, f"{{{PRESENTATION_NS}}}cSld", {"name": self.pdf_reference_slide_name})
        shape_tree = ET.SubElement(common_slide, f"{{{PRESENTATION_NS}}}spTree")

        non_visual_group = ET.SubElement(shape_tree, f"{{{PRESENTATION_NS}}}nvGrpSpPr")
        ET.SubElement(non_visual_group, f"{{{PRESENTATION_NS}}}cNvPr", {"id": "1", "name": ""})
        ET.SubElement(non_visual_group, f"{{{PRESENTATION_NS}}}cNvGrpSpPr")
        ET.SubElement(non_visual_group, f"{{{PRESENTATION_NS}}}nvPr")
        ET.SubElement(shape_tree, f"{{{PRESENTATION_NS}}}grpSpPr")

        picture = ET.SubElement(shape_tree, f"{{{PRESENTATION_NS}}}pic")
        non_visual_picture = ET.SubElement(picture, f"{{{PRESENTATION_NS}}}nvPicPr")
        ET.SubElement(
            non_visual_picture,
            f"{{{PRESENTATION_NS}}}cNvPr",
            {"id": "2", "name": self.pdf_reference_shape_name, "descr": image_description},
        )
        picture_properties = ET.SubElement(non_visual_picture, f"{{{PRESENTATION_NS}}}cNvPicPr")
        ET.SubElement(picture_properties, f"{{{DRAWING_NS}}}picLocks", {"noChangeAspect": "1"})
        ET.SubElement(non_visual_picture, f"{{{PRESENTATION_NS}}}nvPr")

        blip_fill = ET.SubElement(picture, f"{{{PRESENTATION_NS}}}blipFill")
        ET.SubElement(blip_fill, f"{{{DRAWING_NS}}}blip", {f"{{{OFFICE_RELS_NS}}}embed": image_relationship_id})
        stretch = ET.SubElement(blip_fill, f"{{{DRAWING_NS}}}stretch")
        ET.SubElement(stretch, f"{{{DRAWING_NS}}}fillRect")

        shape_properties = ET.SubElement(picture, f"{{{PRESENTATION_NS}}}spPr")
        transform = ET.SubElement(shape_properties, f"{{{DRAWING_NS}}}xfrm")
        ET.SubElement(transform, f"{{{DRAWING_NS}}}off", {"x": "0", "y": "0"})
        ET.SubElement(transform, f"{{{DRAWING_NS}}}ext", {"cx": str(slide_width), "cy": str(slide_height)})
        geometry = ET.SubElement(shape_properties, f"{{{DRAWING_NS}}}prstGeom", {"prst": "rect"})
        ET.SubElement(geometry, f"{{{DRAWING_NS}}}avLst")

        color_map_override = ET.SubElement(root, f"{{{PRESENTATION_NS}}}clrMapOvr")
        ET.SubElement(color_map_override, f"{{{DRAWING_NS}}}masterClrMapping")
        return self._serialize_xml(root)

    def _build_reference_slide_rels_xml(self, *, layout_target: str, image_target: str) -> bytes:
        root = ET.Element(f"{{{PACKAGE_RELS_NS}}}Relationships")
        ET.SubElement(
            root,
            f"{{{PACKAGE_RELS_NS}}}Relationship",
            {"Id": "rId1", "Type": SLIDE_LAYOUT_REL_TYPE, "Target": layout_target},
        )
        ET.SubElement(
            root,
            f"{{{PACKAGE_RELS_NS}}}Relationship",
            {"Id": "rId2", "Type": IMAGE_REL_TYPE, "Target": image_target},
        )
        return self._serialize_xml(root, default_namespace=PACKAGE_RELS_NS)

    @staticmethod
    def _serialize_xml(element: ET.Element, *, default_namespace: str | None = None) -> bytes:
        if default_namespace is not None:
            ET.register_namespace("", default_namespace)
        return ET.tostring(element, encoding="utf-8", xml_declaration=True)
