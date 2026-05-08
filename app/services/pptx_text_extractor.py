from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.models import NormalizedBBox, ParagraphLayout, TextBox, TextLayout, TextSource


class PptxTextExtractor:
    def extract_document(self, pptx_path: Path) -> list[TextLayout]:
        presentation = Presentation(str(pptx_path))
        slide_width = float(presentation.slide_width or 1.0)
        slide_height = float(presentation.slide_height or 1.0)

        layouts: list[TextLayout] = []
        for slide_index, slide in enumerate(presentation.slides):
            paragraphs: list[ParagraphLayout] = []
            lines: list[TextBox] = []
            for paragraph in self._extract_slide_paragraphs(
                slide.shapes,
                slide_index=slide_index,
                slide_width=slide_width,
                slide_height=slide_height,
            ):
                paragraphs.append(paragraph)
                lines.extend(paragraph.lines)

            paragraphs.sort(key=lambda item: (item.bbox[1], item.bbox[0]))
            lines.sort(key=lambda item: (item.bbox[1], item.bbox[0]))
            total_characters = sum(1 for paragraph in paragraphs for character in paragraph.text if not character.isspace())
            layouts.append(
                TextLayout(
                    page_number=slide_index + 1,
                    page_size=(slide_width, slide_height),
                    paragraphs=paragraphs,
                    lines=lines,
                    source=TextSource.NATIVE,
                    total_characters=total_characters,
                    average_confidence=1.0 if paragraphs else 0.0,
                    extracted_with_ocr=False,
                )
            )

        return layouts

    def _extract_slide_paragraphs(
        self,
        shapes,
        *,
        slide_index: int,
        slide_width: float,
        slide_height: float,
        offset_left: int = 0,
        offset_top: int = 0,
    ) -> list[ParagraphLayout]:
        paragraphs: list[ParagraphLayout] = []
        for shape in shapes:
            shape_left = int(getattr(shape, "left", 0) or 0) + offset_left
            shape_top = int(getattr(shape, "top", 0) or 0) + offset_top
            shape_width = int(getattr(shape, "width", 0) or 0)
            shape_height = int(getattr(shape, "height", 0) or 0)
            bbox = self._normalize_bbox(shape_left, shape_top, shape_width, shape_height, slide_width, slide_height)

            if getattr(shape, "has_table", False):
                paragraphs.extend(
                    self._extract_table_paragraphs(
                        shape.table,
                        slide_index=slide_index,
                        bbox=bbox,
                    )
                )
                continue

            if getattr(shape, "has_text_frame", False):
                paragraphs.extend(
                    self._extract_text_frame_paragraphs(
                        shape.text_frame,
                        slide_index=slide_index,
                        bbox=bbox,
                    )
                )

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                try:
                    paragraphs.extend(
                        self._extract_slide_paragraphs(
                            shape.shapes,
                            slide_index=slide_index,
                            slide_width=slide_width,
                            slide_height=slide_height,
                            offset_left=shape_left,
                            offset_top=shape_top,
                        )
                    )
                except Exception:
                    continue

        return paragraphs

    def _extract_table_paragraphs(self, table, *, slide_index: int, bbox: NormalizedBBox) -> list[ParagraphLayout]:
        paragraphs: list[ParagraphLayout] = []
        for row in table.rows:
            for cell in row.cells:
                cell_text = self._clean_text(cell.text)
                if not cell_text:
                    continue
                paragraphs.append(
                    ParagraphLayout(
                        text=cell_text,
                        lines=[
                            TextBox(
                                text=line_text,
                                bbox=bbox,
                                page_number=slide_index + 1,
                                confidence=1.0,
                                source=TextSource.NATIVE,
                            )
                            for line_text in self._split_lines(cell_text)
                        ],
                        bbox=bbox,
                        page_number=slide_index + 1,
                        confidence=1.0,
                        source=TextSource.NATIVE,
                    )
                )
        return paragraphs

    def _extract_text_frame_paragraphs(self, text_frame, *, slide_index: int, bbox: NormalizedBBox) -> list[ParagraphLayout]:
        paragraphs: list[ParagraphLayout] = []
        for paragraph in text_frame.paragraphs:
            paragraph_text = self._clean_text(paragraph.text)
            if not paragraph_text:
                continue
            line_boxes = [
                TextBox(
                    text=line_text,
                    bbox=bbox,
                    page_number=slide_index + 1,
                    confidence=1.0,
                    source=TextSource.NATIVE,
                )
                for line_text in self._split_lines(paragraph_text)
            ]
            paragraphs.append(
                ParagraphLayout(
                    text=paragraph_text,
                    lines=line_boxes,
                    bbox=bbox,
                    page_number=slide_index + 1,
                    confidence=1.0,
                    source=TextSource.NATIVE,
                )
            )
        return paragraphs

    @staticmethod
    def _split_lines(text: str) -> list[str]:
        lines = [part.strip() for part in re.split(r"[\r\n\v]+", text) if part.strip()]
        return lines or [text]

    @staticmethod
    def _clean_text(text: str) -> str:
        cleaned = text.replace("\xa0", " ")
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return cleaned

    @staticmethod
    def _normalize_bbox(
        left: int,
        top: int,
        width: int,
        height: int,
        slide_width: float,
        slide_height: float,
    ) -> NormalizedBBox:
        x0 = max(0.0, min(1.0, float(left) / max(slide_width, 1.0)))
        y0 = max(0.0, min(1.0, float(top) / max(slide_height, 1.0)))
        x1 = max(x0, min(1.0, float(left + max(width, 1)) / max(slide_width, 1.0)))
        y1 = max(y0, min(1.0, float(top + max(height, 1)) / max(slide_height, 1.0)))
        return (x0, y0, x1, y1)
